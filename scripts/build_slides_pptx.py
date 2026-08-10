"""
Build presentation/slide_deck.pptx (native 16:9 PowerPoint) mirroring presentation/slide_deck.md.

Idempotent / re-runnable: overwrites the .pptx in place each run. CPU-only, no network.

Usage:
    python scripts/build_slides_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData

REPO = Path(__file__).resolve().parents[1]
PRES = REPO / "presentation"
IMAGES = PRES / "images"
OUT = PRES / "slide_deck.pptx"

# --- palette (matches dataviz-skill reference categorical order) ---
BLUE = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
GOOD = RGBColor(0x0C, 0xA3, 0x0C)
WARNING = RGBColor(0xFA, 0xB2, 0x19)
CRITICAL = RGBColor(0xD0, 0x3B, 0x3B)
INK = RGBColor(0x20, 0x20, 0x20)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def add_title(slide, text, top=Inches(0.35), size=32, color=BLUE, height=Inches(0.9)):
    box = slide.shapes.add_textbox(Inches(0.5), top, SLIDE_W - Inches(1.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, left=Inches(0.6), top=Inches(1.35), width=Inches(7.6),
                 height=Inches(5.6), size=16, bold_lead=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        if isinstance(item, tuple):
            text, level = item
            p.level = level
        else:
            text = item
        run = p.add_run()
        run.text = f"• {text}" if p.level == 0 else f"– {text}"
        run.font.size = Pt(size - p.level * 1)
        run.font.color.rgb = INK
        p.space_after = Pt(10)
    return box


def add_table(slide, headers, rows, left, top, width, height, header_fill=BLUE,
              font_size=11, header_font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gtable.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(header_font_size)
                r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF2, 0xF2, 0xF0)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = INK
    return gtable


def add_picture_fit(slide, filename, left, top, max_width, max_height):
    path = IMAGES / filename
    if not path.exists():
        raise FileNotFoundError(f"Missing figure: {path}")
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = max_width / max_height
    if aspect > box_aspect:
        w = max_width
        h = Emu(int(max_width / aspect))
    else:
        h = max_height
        w = Emu(int(max_height * aspect))
    left2 = Emu(int(left + (max_width - w) / 2))
    top2 = Emu(int(top + (max_height - h) / 2))
    return slide.shapes.add_picture(str(path), left2, top2, width=w, height=h)


def footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), SLIDE_W - Inches(1.0), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT
    return box


def build():
    prs = new_deck()

    # ============ 1. TITLE SLIDE ============
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    bg.shadow.inherit = False

    box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Invoice Region Detection & Business-Parameter Extraction"
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "From Scanned Invoice to Obligation-Readiness Verdict"
    run2.font.size = Pt(22)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0xCF, 0xE0, 0xF2)

    box2 = s.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.2))
    tf2 = box2.text_frame
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = "Rolando · Diana · Jordan · Damir · Hessam"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf2.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Deep Learning II — Group Project"
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0xCF, 0xE0, 0xF2)

    set_notes(s,
        "Welcome everyone. Today we're presenting an end-to-end pipeline that takes a scanned "
        "invoice image and decides whether it's ready to become a structured digital obligation "
        "record - with a transparent, user-configurable verdict, not a black box. Five of us, "
        "five stages, one shared interface contract.")

    # ============ 2. THE BUSINESS PROBLEM ============
    s = blank_slide(prs)
    add_title(s, "The Business Problem")
    add_bullets(s, [
        "A business receives thousands of invoices as images/scans",
        "Before an invoice can feed an automated finance workflow, someone must check:",
        ("Is it signed or stamped?", 1),
        ("Does it cite a PO / contract / work order?", 1),
        ("Does it state clear payment terms and a valid date?", 1),
        "Manual triage doesn't scale — slow, inconsistent, delays payment",
        "Goal: read an invoice image → structured JSON + a transparent Ready / Not-ready verdict",
    ])
    set_notes(s,
        "This is the Pistac.io framing: a digital obligation record is a machine-readable summary "
        "of what an invoice obligates the business to do, and by when, that downstream systems can "
        "act on without a human reading the source image first. At real volume, manual review is "
        "the bottleneck we're targeting.")

    # ============ 3. WHY VISION + OCR ============
    s = blank_slide(prs)
    add_title(s, "Why This Is a Vision + OCR Problem")
    add_bullets(s, [
        "Not just \"read the text\" — first you must know where to look",
        "Detect visual marks (stamp/signature) and business regions before running OCR",
        "OCR only the relevant crops → faster, more accurate than OCR-everything",
        "Then apply rule-based extraction to the OCR'd text for structured fields",
        "Finally, judge the extracted signals against a user-defined policy",
    ])
    set_notes(s,
        "We split the problem into detection (where), reading (what), and judgment (is it "
        "enough). That separation is what let five people work in parallel on independent, "
        "testable stages.")

    # ============ 4. PIPELINE OVERVIEW ============
    s = blank_slide(prs)
    add_title(s, "Pipeline Overview")
    stages = [
        ("01 Rolando", "Data ingestion\ninvoice_manifest.csv", ACCENT),
        ("02 Diana", "Stamp/Signature\nYOLOv8n", ORANGE),
        ("03 Jordan", "Region Detection\nYOLOv8n", AQUA),
        ("04 Damir", "OCR + Params/Terms\nEasyOCR + rules", WARNING),
        ("05 Hessam", "Integration +\nVerdict Engine +\nStreamlit App", BLUE),
    ]
    box_w = Inches(2.15)
    box_h = Inches(1.6)
    gap = Inches(0.28)
    total_w = box_w * 5 + gap * 4
    start_x = int((SLIDE_W - total_w) / 2)
    y = Inches(2.6)
    xs = []
    for i, (title, sub, color) in enumerate(stages):
        x = Emu(int(start_x + i * (box_w + gap)))
        xs.append(x)
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = WHITE
        shp.line.width = Pt(1)
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sub
        run2.font.size = Pt(10.5)
        run2.font.color.rgb = WHITE
    for i in range(4):
        ax = Emu(int(xs[i] + box_w))
        ay = Emu(int(y + box_h / 2 - Inches(0.08)))
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, gap, Inches(0.16))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()
        arrow.shadow.inherit = False
    out_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.4), Inches(4.7), Inches(6.5), Inches(0.9)
    )
    out_box.fill.solid()
    out_box.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF0)
    out_box.line.color.rgb = MUTED
    tf = out_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Final JSON per invoice + Ready/Not-ready verdict"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = INK
    connector = s.shapes.add_connector(2, Inches(6.65), Inches(4.2), Inches(6.65), Inches(4.7))
    connector.line.color.rgb = MUTED
    set_notes(s,
        "Each notebook publishes its outputs to Google Drive; Hessam's integration stage consumes "
        "all four upstream outputs and produces one JSON per invoice plus the verdict. "
        "document_id is the join key across every CSV.")

    # ============ 5. DATASETS / DOMAIN GAP ============
    s = blank_slide(prs)
    add_title(s, "Datasets — and an Honest Domain Gap")
    add_table(
        s,
        ["Dataset", "Contents", "Used by"],
        [
            ["Invoice batches", "~5,201 unique real scans; 750 sampled; batch_1 has real "
             "annotations", "Rolando, all"],
            ["OCR Dataset of Multi-type Documents", "973 receipts, 100% annotated, 52,331 real "
             "boxes", "Jordan, Damir"],
            ["SignverOD", "2,765 document images, signature boxes", "Diana"],
            ["StaVer", "400 scans, stamp ground-truth masks", "Diana"],
        ],
        Inches(0.6), Inches(1.3), Inches(12.1), Inches(2.6),
        font_size=13, header_font_size=13,
    )
    add_bullets(s, [
        "None of these are invoices. Detectors train on the best real labelled data, then get applied to invoices",
        "On invoices: we report detection counts, not accuracy — there's no invoice-level ground truth",
    ], top=Inches(4.2), height=Inches(1.4), size=16)
    set_notes(s,
        "Say this part out loud - it's a strength, not a weakness, to be upfront about it. Our "
        "metrics are real numbers on each dataset's own held-out split. Applying to invoices is "
        "an honest, disclosed domain-transfer step.")

    # ============ 6. ROLANDO ============
    s = blank_slide(prs)
    add_title(s, "Rolando — Data Ingestion & the Manifest")
    add_bullets(s, [
        "No model — this is data engineering, and the whole pipeline's honesty starts here",
        "Builds invoice_manifest.csv: document_id, image_path, width, height, split, has_ground_truth",
        "Annotation-aware resampling: naive sampling → 26.3% GT coverage → resampled to 100% "
        "(750 rows, split 525/120/105)",
        "Caught real data leakage: batch_3 secretly duplicates batches 1-2 → true unique count "
        "= 5,201, not 8,181 → duplicates excluded",
    ], width=Inches(6.6))
    add_picture_fit(s, "sample_invoice_grid.png", Inches(7.4), Inches(1.4), Inches(5.4), Inches(5.4))
    set_notes(s,
        "Garbage in, garbage out. If the manifest is wrong, every downstream metric is wrong. The "
        "annotation-aware resampling was a real trade-off: full labels vs. cross-batch visual "
        "variety. Catching the batch_3 duplication before it leaked into train/test splits was one "
        "of the more important engineering wins in this project.")

    # ============ 7. DIANA ============
    s = blank_slide(prs)
    add_title(s, "Diana — Stamp & Signature Detection")
    add_bullets(s, [
        "Model: YOLOv8n, single 2-class detector (stamp, signature — never merged)",
        "SignverOD: normalized boxes → pixels; only category 1 (signature) kept",
        "StaVer: ships no boxes, only masks → boxes derived via connectedComponentsWithStats",
        "Budget-optimized run: ≤50 epochs, imgsz 640 (reduced after GPU exhaustion at epoch 86)",
        "Invoice inference: 0 of 750 invoices have any detection — corpus is clean, unsigned "
        "digital templates (honest domain-gap result, not a model failure)",
    ], width=Inches(6.6), size=15)
    add_table(
        s, ["Class", "Precision", "Recall", "Mean IoU"],
        [["stamp", "0.903", "0.875", "0.822"], ["signature", "0.894", "0.638", "0.815"]],
        Inches(0.6), Inches(4.9), Inches(6.4), Inches(1.4), header_fill=ORANGE,
    )
    add_picture_fit(s, "stamp_signature_detection_examples.png", Inches(7.3), Inches(1.35),
                     Inches(5.5), Inches(5.6))
    set_notes(s,
        "Report per-class, never a blended score - a 'must be signed' business rule cares about "
        "the signature number specifically. Signature recall (0.638) is lower than precision "
        "(0.894) - the model is conservative, it misses some real signatures (382 false "
        "negatives) rather than over-calling them. The mask-to-box derivation for StaVer was the "
        "trickiest engineering piece of this notebook - know it cold for questions.")

    # ============ 8. JORDAN ============
    s = blank_slide(prs)
    add_title(s, "Jordan — Region Detection & IoU")
    add_bullets(s, [
        "Model: YOLOv8n, 5-class detector: company, date, address, total, other_text",
        "OCR Dataset gives text boxes with no class + entity values with no coordinates",
        "Join by fuzzy-matching entity text to box text (rapidfuzz, threshold 88)",
        "Challenge: other_text massively outnumbers the four field classes → report per-class, "
        "never a single mAP",
        "Invoice inference: 750 of 750 invoices carry a region detection (domain shift, counts only)",
    ], width=Inches(6.6), size=15)
    add_table(
        s, ["Class", "Precision", "Recall", "Mean IoU"],
        [
            ["company", "0.897", "0.819", "0.894"],
            ["date", "0.756", "0.855", "0.812"],
            ["address", "0.880", "0.883", "0.896"],
            ["total", "0.763", "0.783", "0.887"],
            ["other_text", "0.911", "0.967", "0.877"],
            ["Macro mean IoU", "—", "—", "0.873"],
        ],
        Inches(0.6), Inches(4.5), Inches(6.4), Inches(2.7), header_fill=AQUA, font_size=10.5,
    )
    add_picture_fit(s, "region_detection_examples.png", Inches(7.3), Inches(1.35),
                     Inches(5.5), Inches(5.6))
    set_notes(s,
        "Jordan's date/total/company boxes get a second life in the app: crop + OCR just that "
        "region to localize and read key fields - this feeds the verdict engine's date rule. "
        "date has the lowest precision (0.756) of the four named classes - dates are short, "
        "numeric, and easy to confuse with other short numeric text on a receipt.")

    # ============ 9. DAMIR ============
    s = blank_slide(prs)
    add_title(s, "Damir — OCR, Parameters & Terms")
    add_bullets(s, [
        "OCR: EasyOCR (deep-learning, GPU) — vs Tesseract/PaddleOCR: stronger accuracy, simple setup",
        "Rule-based extraction (parameter_checker.py, terms_extraction.py) — transparent, "
        "user-extensible via config",
        "Two eval sets, always reported together, never one headline number:",
        ("Primary (98 receipts, real GT): CER 0.2152 / WER 0.5423", 1),
        ("Secondary (120 clean invoices): CER 0.0002 / WER 0.0015", 1),
        "Payment terms nearly absent from corpus: 99.6% have a date, only 0.4% have day-based terms",
    ], width=Inches(8.6), size=16)
    add_table(
        s, ["Metric", "Receipts (primary)", "Invoices (secondary)"],
        [
            ["CER (mean)", "0.2152", "0.0002"],
            ["WER (mean)", "0.5423", "0.0015"],
            ["Docs scored", "98", "120 / 750 with GT"],
        ],
        Inches(0.6), Inches(5.5), Inches(11.8), Inches(1.5), header_fill=WARNING, font_size=12,
    )
    set_notes(s,
        "Always name the denominator when reporting the secondary number - it's honest and it's a "
        "strength to volunteer this, not something a professor has to extract with a follow-up "
        "question. Required-field presence: PO Reference 54.7%, Order Number 19.9% - the two that "
        "actually gate readiness.")

    # ============ 10. HESSAM ============
    s = blank_slide(prs)
    add_title(s, "Hessam — Integration & the Final JSON")
    add_bullets(s, [
        "Merges Diana + Jordan + Damir outputs into one JSON per invoice (document_id is the join key)",
        "Schema fixed by model_interface_contract.md — visual elements, detected regions, "
        "required parameters, payment context, terms & conditions, readiness, model metrics",
        "Integration nuance: Damir's parameter/terms CSVs are keyed to the receipt dataset, not "
        "the 750 invoices → invoice-level signals are derived at integration",
        "750 fused per-invoice JSON records produced",
    ], width=Inches(8.6), size=16)
    add_table(
        s, ["Metric", "Value"],
        [
            ["Invoices in manifest", "750"],
            ["Final-JSON records produced", "750"],
            ["Region boxes on invoices (other_text)", "54,009"],
        ],
        Inches(0.6), Inches(5.0), Inches(11.8), Inches(1.7), header_fill=BLUE, font_size=13,
    )
    set_notes(s,
        "This is a subtle but important point: we didn't fork Damir's logic, we re-ran his exact "
        "shared functions against a different text source at integration time - so results stay "
        "comparable and auditable back to one implementation.")

    # ============ 11. VERDICT ENGINE ============
    s = blank_slide(prs)
    add_title(s, "The Verdict Engine — User-Configurable, Fail-Closed")
    add_bullets(s, [
        "The user builds a policy from four toggleable rule types before judging any invoice:",
        ("Visual mark — stamp OR/AND signature (or either alone)", 1),
        ("Reference number — any-of PO / Order / Contract / Work Order", 1),
        ("Invoice date range", 1),
        ("Payment terms — e.g. > 30 days", 1),
        "Verdict = AND of every enabled rule",
        "Fail-closed: missing evidence (unknown) counts as fail — never a silent pass",
        "Pure logic, no I/O — independently unit-tested (src/verdict_engine.py, 9/9 passing)",
    ])
    set_notes(s,
        "This is the headline feature. A compliance officer, an AP clerk, and an auditor may "
        "reasonably want different thresholds - so the policy is exposed as configuration, not "
        "buried in code. And a readiness gate must never pass what it cannot confirm - that's why "
        "we chose fail-closed as the default, and we show it transparently in the breakdown rather "
        "than hiding the 'unknown' case.")

    # ============ 12. READINESS SPREAD ============
    s = blank_slide(prs)
    add_title(s, "The Readiness Spread — Same 750 Invoices, Two Views")
    add_picture_fit(s, "readiness_by_policy.png", Inches(2.4), Inches(1.15), Inches(8.5), Inches(4.6))
    add_bullets(s, [
        "Strict policies (fail-closed): Lenient 750/750 (100%), Default ~0/750 (~0%), Strict 0/750 (0%)",
        "Default ~0% = corpus carries no PO/contract references (earlier 63.3% was a matching artifact, corrected)",
        "Graded completeness: Ready 315/750 (42.0%), Needs review 93 (12.4%), Not ready 342 (45.6%), mean 76.3",
    ], top=Inches(5.85), height=Inches(1.5), size=14)
    set_notes(s,
        "The verdict engine is a real, configurable system - same 750 invoices, different honest "
        "outcomes. The strict Default is ~0% because the corpus carries no PO/contract references "
        "(an earlier 63.3% there was a loose-matching artifact, since corrected); Strict is 0% "
        "because there are no stamps/signatures - fail-closed working as intended. Because the "
        "strict gate can't differentiate this corpus, we add a standardized graded completeness "
        "score (Ready 42%, Needs review 12%, Not ready 46%) that yields a meaningful, auditable spread.")

    # ============ 13. LIVE DEMO ============
    s = blank_slide(prs)
    add_title(s, "Live Demo — the Streamlit App")
    add_bullets(s, [
        "Live Demo — upload or pick one invoice → annotated boxes → verdict card with "
        "per-rule breakdown → JSON / report download",
        "Batch Gallery — apply the policy across all 750 invoices → \"N of 750 pass\" → "
        "filter, drill down, export the passing set",
        "Model Report — each member's real metrics + local-CPU-vs-Colab-GPU comparison charts",
        "Hybrid execution: gallery reads pre-computed results; a new upload runs the live pipeline",
    ])
    set_notes(s,
        "[SHOW: switch policy presets live - Default / Strict / Lenient - and watch the verdict "
        "and the batch roll-up count change in real time.] This is the moment that makes "
        "'configurable' concrete instead of just a claim.")

    # ============ 14. RESULTS ============
    s = blank_slide(prs)
    add_title(s, "Results")
    add_table(
        s, ["Stage", "Headline metric", "Value"],
        [
            ["Diana (stamp/signature)", "Precision / Recall / Mean IoU per class",
             "stamp 0.903/0.875/0.822 — signature 0.894/0.638/0.815"],
            ["Jordan (regions)", "Per-class Mean IoU, macro mean IoU",
             "0.812–0.896 per class — macro 0.873"],
            ["Damir (OCR)", "CER / WER primary / CER secondary",
             "0.2152 / 0.5423 (receipts) — 0.0002 (invoices)"],
            ["Hessam (integration)", "Readiness: strict / graded", "Default ~0%, Lenient 100%, Strict 0% / Ready 42.0% (315/750)"],
        ],
        Inches(0.6), Inches(1.5), Inches(12.1), Inches(3.2), header_fill=BLUE, font_size=13,
    )
    add_picture_fit(s, "metrics_per_class.png", Inches(1.4), Inches(4.9), Inches(10.5), Inches(2.2))
    set_notes(s,
        "These are the final numbers - traceable to each metrics JSON or the final_pipeline_report.md. "
        "State the strict readiness spread (Lenient 100% / Default ~0% / Strict 0%) and the graded "
        "completeness score (Ready 42%) if asked; the earlier Default 63.3% was a loose-matching "
        "artifact, since corrected, and Strict=0% is the fail-closed design on an unsigned corpus.")

    # ============ 15. LIMITATIONS ============
    s = blank_slide(prs)
    add_title(s, "Limitations & Honest Caveats")
    add_bullets(s, [
        "Domain gap: detectors train on documents/receipts, get applied to invoices — counts, "
        "not accuracy, on invoices",
        "GT coverage: annotation CSVs exist only for batch_1 — resampling lifted coverage "
        "26.3% → 100%, trading cross-batch variety for full labels",
        "Data leakage caught: batch_3 duplicated batches 1-2 — true unique invoices = 5,201, "
        "not 8,181",
        "Integration re-derivation: Damir's receipt-keyed CSVs → invoice signals derived at "
        "integration; batch-path coverage is bounded by available OCR/annotation text",
        "Fail-closed by design: \"Not-ready\" includes both real failures and unconfirmed checks "
        "— always shown separately in the breakdown",
        "Compute budget: Diana's model trained at a reduced budget (≤50 epochs, imgsz 640) after "
        "a GPU-budget exhaustion at epoch 86",
    ], size=15)
    set_notes(s,
        "Every one of these is stated plainly because owning a limitation is a stronger position "
        "than hiding it - and because it's true. Rehearse being able to explain why each "
        "trade-off was made, not just that it exists.")

    # ============ 16. CONCLUSION ============
    s = blank_slide(prs)
    add_title(s, "Conclusion & Future Work")
    add_bullets(s, [
        "A business-facing readiness question decomposes into independently trainable/evaluable "
        "stages, integrated through a fixed schema, judged by a transparent, configurable, "
        "fail-closed engine",
        "Every stage evaluated on real, labelled data with real metrics; domain gap disclosed, "
        "not hidden",
        "Future work:",
        ("Collect invoice-native ground truth for stamp/signature/region boxes", 1),
        ("Extend annotation-aware sampling beyond batch_1", 1),
        ("A learned OCR-quality gate to distinguish \"absent\" from \"OCR failed\"", 1),
        ("Expose confidence/fuzzy-match thresholds as live sidebar controls", 1),
    ])
    set_notes(s,
        "Close on the honesty theme: the biggest lever for improving this system isn't a bigger "
        "model, it's more invoice-native ground truth. That's a data problem, not a modeling "
        "problem, and we want to be clear about that distinction.")

    # ============ 17. THANK YOU ============
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    bg.shadow.inherit = False
    box = s.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Thank You — Questions?"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    add_bullets(s, [
        "Rolando — Data Ingestion & Manifest",
        "Diana — Stamp & Signature Detection",
        "Jordan — Region Detection & IoU",
        "Damir — OCR, Parameters & Terms",
        "Hessam — Integration, Verdict Engine & Streamlit App",
    ], top=Inches(2.4), size=18)
    for shape in s.shapes:
        if shape.has_text_frame and shape is not box:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.color.rgb == INK:
                        r.font.color.rgb = WHITE
    set_notes(s,
        "Open the floor. Each of us can speak in depth to our own stage - route detection/model "
        "questions to Diana or Jordan, OCR/extraction questions to Damir, data-quality questions "
        "to Rolando, and integration/app/verdict-engine questions to Hessam.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")
    return prs


def verify():
    prs = Presentation(str(OUT))
    n_slides = len(prs.slides._sldIdLst)
    print(f"Reload check: {OUT.name} -> {n_slides} slides")
    return n_slides


if __name__ == "__main__":
    build()
    verify()
